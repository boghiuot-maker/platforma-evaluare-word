from pptx import Presentation
from pathlib import Path
from fpdf import FPDF
from openpyxl import Workbook

def emu_to_cm(emu):
    return emu / 360000.0

def text_of_title(slide):
    try:
        return slide.shapes.title.text.strip()
    except Exception:
        return ""

def count_paragraphs_in_slide(slide):
    for shp in slide.shapes:
        try:
            if shp.has_text_frame:
                pats = [p for p in shp.text_frame.paragraphs if p.text.strip()]
                if pats:
                    return len(pats), [p.text.strip() for p in pats]
        except Exception:
            continue
    return 0, []

def find_pictures(slide):
    pics = []
    for shp in slide.shapes:
        try:
            if getattr(shp, "shape_type", None) == 13:
                pics.append(shp)
        except Exception:
            pass
    return pics

def shape_xml_contains(shp, token):
    try:
        return token.lower() in shp._element.xml.lower()
    except Exception:
        return False

def slide_xml_contains(slide, token):
    try:
        return token.lower() in slide._element.xml.lower()
    except Exception:
        return False

def evaluate_pptx(path:Path):
    prs = Presentation(str(path))
    info = {"slides":len(prs.slides),"titles":[], "texts":[],"pics":0,"notes":""}
    for i, slide in enumerate(prs.slides, start=1):
        info["titles"].append(text_of_title(slide))
        # text snippet
        cnt, texts = count_paragraphs_in_slide(slide)
        info["texts"].append(texts[0] if texts else "")
        pics = find_pictures(slide)
        info["pics"] += len(pics)
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                info["notes"] += notes + " ; "
        except Exception:
            pass
    return info

def evaluate_student_package(folder:Path):
    results = {}
    for i in range(1,7):
        key=f"app{i}"
        results[key] = {"pptx":None,"report":None,"capture":None,"export_pdf":None}
        for f in folder.iterdir():
            if f.name.startswith(f"app{i}_"):
                results[key]["pptx"] = f
                try:
                    results[key]["report"] = evaluate_pptx(f)
                except Exception as e:
                    results[key]["report"] = {"error":str(e)}
            if f.name.startswith(f"cap{i}_"):
                results[key]["capture"] = f
        # export pdf may be stored as 'rezolvare.pdf' at folder root
        ep = folder / "rezolvare.pdf"
        if ep.exists():
            results[key]["export_pdf"] = ep
    return results

def generate_report_files(folder:Path, name, klass, date, results, reports_dir:Path):
    reports_dir.mkdir(parents=True, exist_ok=True)
    base = reports_dir / f"{name.replace(' ','_')}_{klass}_{folder.name}"
    txtp = base.with_suffix(".txt")
    pdfp = base.with_suffix(".pdf")
    xlsxp = base.with_suffix(".xlsx")

    lines = [f"Evaluare PowerPoint pentru: {name} | Clasa: {klass} | Data: {date}", ""]
    total_ok = 0
    # run checks for each requirement
    # Map checks to slides: slides 1..4 used etc
    # We'll create detailed per requirement messages
    checks = []
    # Requirement 1: exact 4 slides
    any_ppt = None
    for i in range(1,7):
        if results[f"app{i}"]["pptx"]:
            any_ppt = results[f"app{i}"]["pptx"]
            break
    if not any_ppt:
        lines.append("No PPTX files uploaded for evaluation.")
        txtp.write_text("\\n".join(lines), encoding='utf-8')
        # still create empty pdf/xlsx
        pdf = FPDF()
        pdf.add_page(); pdf.set_font("Arial", size=12); pdf.multi_cell(0,6,"No files")
        pdf.output(str(pdfp))
        wb = Workbook(); wb.save(str(xlsxp))
        return pdfp, xlsxp

    # We'll evaluate using the first uploaded PPTX as the student's presentation
    ppt_path = any_ppt
    rpt = results["app1"]["report"] if results["app1"]["report"] else (results["app2"]["report"] or results["app3"]["report"])
    # load presentation to inspect slides in detail
    from pptx import Presentation as Pres
    pres = Pres(str(ppt_path))
    slide_count = len(pres.slides)
    # Req1
    r1 = ("OK" if slide_count==4 else f"NU (found {slide_count} slides)")
    checks.append(("Cerința 1", r1))
    if r1.startswith("OK"): total_ok += 1
    # Req2 title slide1
    t1 = ""
    try:
        t = pres.slides[0].shapes.title.text.strip() if pres.slides[0].shapes.title else ""
    except Exception:
        t = ""
    r2 = ("OK" if t=="Evaluare competențe PowerPoint" else f'NU (found: \"{t}\")')
    checks.append(("Cerința 2", r2))
    if r2.startswith("OK"): total_ok += 1
    # Req3 slide3 title
    t3 = ""
    try:
        t3 = pres.slides[2].shapes.title.text.strip() if pres.slides[2].shapes.title else ""
    except Exception:
        t3 = ""
    r3 = ("OK" if t3=="Structură" else f'NU (found: \"{t3}\")')
    checks.append(("Cerința 3", r3))
    if r3.startswith("OK"): total_ok += 1
    # Req4 list on slide2 with 3 bullets
    p_count = 0; bullets = []
    try:
        slide2 = pres.slides[1]
        for shp in slide2.shapes:
            if shp.has_text_frame:
                paras = [p.text.strip() for p in shp.text_frame.paragraphs if p.text.strip()]
                if paras:
                    p_count = len(paras); bullets = paras; break
    except Exception:
        p_count = 0
    r4 = ("OK" if p_count==3 else f"NU (found {p_count} items)")
    checks.append(("Cerința 4", r4))
    if r4.startswith("OK"): total_ok += 1
    # Req5 first bullet text equals Introducere
    r5 = ("OK" if bullets and bullets[0]=="Introducere" else f'NU (found: \"{bullets[0] if bullets else \"\"}\")')
    checks.append(("Cerința 5", r5))
    if r5.startswith("OK"): total_ok += 1
    # Req6 bullet type (best-effort: check for bullet char in xml)
    try:
        slide2_xml = pres.slides[1]._element.xml.lower()
        r6 = ("OK" if "<a:buChar" in slide2_xml or "<a:buFont" in slide2_xml or "lvl" in slide2_xml else "NU (bullet type not standard detected)")
    except Exception:
        r6 = "NU (could not read xml)"
    checks.append(("Cerința 6", r6))
    if r6.startswith("OK"): total_ok += 1
    # Req7 picture on slide3
    pics3 = sum(1 for shp in pres.slides[2].shapes if getattr(shp,'shape_type',None)==13)
    r7 = ("OK" if pics3>=1 else "NU (no picture)")
    checks.append(("Cerința 7", r7))
    if r7.startswith("OK"): total_ok += 1
    # Req8 width between 6.5 and 7.5 cm for first picture on slide3
    r8 = "NU (no picture)"
    if pics3>=1:
        for shp in pres.slides[2].shapes:
            try:
                if getattr(shp,'shape_type',None)==13:
                    w_cm = shp.width/360000.0
                    if 6.5 <= w_cm <= 7.5:
                        r8 = "OK"
                    else:
                        r8 = f"NU (width {w_cm:.2f} cm)"
                    break
            except Exception:
                r8 = "NU (error reading size)"
    checks.append(("Cerința 8", r8))
    if r8.startswith("OK"): total_ok += 1
    # Req9 alt text 'Imagine exemplu' on picture
    r9 = "NU (no picture)"
    found_alt = False
    for shp in pres.slides[2].shapes:
        try:
            if getattr(shp,'shape_type',None)==13:
                xml = shp._element.xml.lower()
                if 'descr="' in xml and 'imagine exemplu' in xml:
                    found_alt = True; break
        except Exception:
            pass
    r9 = ("OK" if found_alt else "NU (alt text not found)")
    checks.append(("Cerința 9", r9))
    if r9.startswith("OK"): total_ok += 1
    # Req10 table on slide4 with 3 cols and 2 rows
    r10 = "NU (no table)"
    try:
        for shp in pres.slides[3].shapes:
            if getattr(shp,'has_table',False):
                tbl = shp.table
                if tbl.columns==3 and tbl.rows==2:
                    r10 = "OK"
                else:
                    r10 = f"NU (found {tbl.rows}x{tbl.columns})"
                break
    except Exception:
        r10 = "NU (error)"
    checks.append(("Cerința 10", r10))
    if r10.startswith("OK"): total_ok += 1
    # Req11 first cell text "Total"
    r11 = "NU (no table)"
    try:
        for shp in pres.slides[3].shapes:
            if getattr(shp,'has_table',False):
                tbl = shp.table
                cell = tbl.cell(0,0).text.strip()
                r11 = ("OK" if cell=="Total" else f'NU (found: \"{cell}\")')
                break
    except Exception:
        r11 = "NU (error)"
    checks.append(("Cerința 11", r11))
    if r11.startswith("OK"): total_ok += 1
    # Req12 chart presence on slide4
    r12 = "NU (no chart)"
    try:
        has_chart = any(getattr(shp,'has_chart',False) for shp in pres.slides[3].shapes)
        r12 = ("OK" if has_chart else "NU (no chart)")
    except Exception:
        r12 = "NU (error)"
    checks.append(("Cerința 12", r12))
    if r12.startswith("OK"): total_ok += 1
    # Req13 transition fade on slide1 (best-effort xml search)
    try:
        s1xml = pres.slides[0]._element.xml.lower()
        r13 = ("OK" if 'transition' in s1xml and 'fade' in s1xml else "NU (fade transition not found)")
    except Exception:
        r13 = "NU (error)"
    checks.append(("Cerința 13", r13))
    if r13.startswith("OK"): total_ok += 1
    # Req14 duration 1 second (1000ms) best-effort search
    try:
        r14 = ("OK" if 'dur' in pres.slides[0]._element.xml.lower() and '1000' in pres.slides[0]._element.xml.lower() else "NU (duration not 1000ms or not found)")
    except Exception:
        r14 = "NU (error)"
    checks.append(("Cerința 14", r14))
    if r14.startswith("OK"): total_ok += 1
    # Req15 animation wipe on slide2
    try:
        s2xml = pres.slides[1]._element.xml.lower()
        r15 = ("OK" if 'wipe' in s2xml else "NU (wipe animation not found)")
    except Exception:
        r15 = "NU (error)"
    checks.append(("Cerința 15", r15))
    if r15.startswith("OK"): total_ok += 1
    # Req16 hyperlink to icdl.org on slide1
    found_hlink = False
    try:
        for shp in pres.slides[0].shapes:
            try:
                if shp.has_text_frame:
                    for p in shp.text_frame.paragraphs:
                        for r in p.runs:
                            hl = getattr(r, 'hyperlink', None)
                            if hl and getattr(hl, 'address', None):
                                if 'icdl.org' in hl.address:
                                    found_hlink = True; break
            except Exception:
                pass
    except Exception:
        pass
    r16 = ("OK" if found_hlink else "NU (hyperlink not found)")
    checks.append(("Cerința 16", r16))
    if r16.startswith("OK"): total_ok += 1
    # Req17 notes on slide3 contain text
    try:
        notes = pres.slides[2].notes_slide.notes_text_frame.text.strip() if pres.slides[2].has_notes_slide else ""
        r17 = ("OK" if "Notă pentru evaluator" in notes else f'NU (found: \"{notes}\")')
    except Exception:
        r17 = "NU (error)"
    checks.append(("Cerința 17", r17))
    if r17.startswith("OK"): total_ok += 1
    # Req18 export PDF named rezolvare.pdf present in folder
    exp = folder = Path(ppt_path).parent
    pdff = exp / "rezolvare.pdf"
    r18 = ("OK" if pdff.exists() else "NU (rezolvare.pdf not uploaded)")
    checks.append(("Cerința 18", r18))
    if r18.startswith("OK"): total_ok += 1

    # compile report
    for c in checks:
        lines.append(f"{c[0]}: {c[1]}")
    lines.append("")
    lines.append(f"Punctaj automat (items OK): {total_ok} / 18 (demo metric)")

    txtp.write_text("\\n".join(lines), encoding='utf-8')

    # PDF generation with basic font
    pdf = FPDF()
    pdf.add_page()
    try:
        pdf.add_font("DejaVu","", "static/fonts/DejaVuSans.ttf", uni=True)
        pdf.set_font("DejaVu", size=12)
    except Exception:
        pdf.set_font("Arial", size=12)
    for ln in lines:
        pdf.multi_cell(0,6,ln)
    pdf.output(str(pdfp))

    # XLSX
    wb = Workbook()
    ws = wb.active
    ws.title = "Detalii"
    ws.append(["Cerința","Rezultat"])
    for c in checks:
        ws.append([c[0], c[1]])
    wb.save(str(xlsxp))

    return pdfp, xlsxp
